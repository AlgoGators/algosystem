from pptx import Presentation
from pptx.util import Inches
import pandas as pd
import os
from datetime import datetime
import argparse
import sys

'''
# Load your template presentation
prs = Presentation("template.pptx")

# Select the slide you want to modify (e.g., first slide)
slide = prs.slides[0]

# Add an image to that slide
img_path = "example.png"
left = Inches(1)
top = Inches(2)
height = Inches(3)
slide.shapes.add_picture(img_path, left, top, height=height)

# Optionally add text to a placeholder
slide.placeholders[0].text = "Updated Title from Python"

# Save as a new PowerPoint file
prs.save("output_with_image.pptx")
'''

def export_backtest_to_csv(results, output_dir="backtest_exports", prefix="backtest"):
    """
    Export backtest results including metrics and graph data to CSV files.
    
    Parameters
    ----------
    results : dict
        Dictionary containing backtest results from engine.run()
    output_dir : str, default="backtest_exports"
        Directory to save the CSV files
    prefix : str, default="backtest"
        Prefix for the exported files
    
    Returns
    -------
    dict
        Dictionary containing paths to all exported files
    """
    # Create output directory if it doesn't exist
    os.makedirs(output_dir, exist_ok=True)
    
    # Generate timestamp for unique filenames
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    
    exported_files = {}
    
    # 1. Export Summary Metrics
    summary_data = {
        'Metric': ['Initial Capital', 'Final Capital', 'Total Return', 
                   'Start Date', 'End Date'],
        'Value': [
            results.get('initial_capital', 'N/A'),
            results.get('final_capital', 'N/A'),
            f"{results.get('returns', 0) * 100:.2f}%",
            results.get('start_date', 'N/A'),
            results.get('end_date', 'N/A')
        ]
    }
    
    # Add all calculated metrics
    metrics_data = results.get('metrics', {})
    for key, value in metrics_data.items():
        if not key.endswith('_error'):
            metric_name = key.replace('_', ' ').title()
            if isinstance(value, (int, float)):
                if 'return' in key.lower() or 'ratio' in key.lower():
                    formatted_value = f"{value * 100:.2f}%" if value < 1 else f"{value:.2f}"
                else:
                    formatted_value = f"{value:.4f}"
            else:
                formatted_value = str(value)
            
            summary_data['Metric'].append(metric_name)
            summary_data['Value'].append(formatted_value)
    
    summary_df = pd.DataFrame(summary_data)
    summary_file = os.path.join(output_dir, f"{prefix}_summary_{timestamp}.csv")
    summary_df.to_csv(summary_file, index=False)
    exported_files['summary'] = summary_file
    
    # 2. Export Time Series Data for Graphs
    plots_data = results.get('plots', {})
    
    time_series_dfs = {}
    for key, series_data in plots_data.items():
        if isinstance(series_data, pd.Series):
            col_name = key.replace('_', ' ').title()
            time_series_dfs[col_name] = series_data
    
    if time_series_dfs:
        combined_ts = pd.DataFrame(time_series_dfs)
        combined_ts.index.name = 'Date'
        
        timeseries_file = os.path.join(output_dir, f"{prefix}_timeseries_{timestamp}.csv")
        combined_ts.to_csv(timeseries_file)
        exported_files['timeseries'] = timeseries_file
        
        # Export specialized chart data
        _export_chart_specific_data(combined_ts, output_dir, prefix, timestamp, exported_files)
    
    # 3. Export Raw Equity Series
    if 'equity' in results and isinstance(results['equity'], pd.Series):
        equity_series = pd.DataFrame({
            'Date': results['equity'].index,
            'Portfolio Value': results['equity'].values
        })
        raw_equity_file = os.path.join(output_dir, f"{prefix}_portfolio_value_{timestamp}.csv")
        equity_series.to_csv(raw_equity_file, index=False)
        exported_files['portfolio_value'] = raw_equity_file
    
    # 4. Create metadata file
    _create_metadata_file(results, exported_files, output_dir, prefix, timestamp)
    
    charts = create_charts_from_exports(
        output_dir="backtest_exports",
        prefix="backtest",
        charts_dir="charts",
        pptx_template="template.pptx",   # or None
        pptx_out="backtest_charts.pptx"  # or None to skip pptx creation
    )

    return exported_files


def _export_chart_specific_data(combined_ts, output_dir, prefix, timestamp, exported_files):
    """Helper function to export chart-specific data"""
    # Equity Curve Data
    equity_data = pd.DataFrame()
    for col in ['Equity Curve', 'Benchmark Equity Curve', 'Relative Performance']:
        if col in combined_ts.columns:
            equity_data[col] = combined_ts[col]
    
    if not equity_data.empty:
        equity_file = os.path.join(output_dir, f"{prefix}_equity_curve_{timestamp}.csv")
        equity_data.to_csv(equity_file)
        exported_files['equity_curve'] = equity_file
    
    # Drawdown Data
    drawdown_data = pd.DataFrame()
    for col in ['Drawdown Series', 'Benchmark Drawdown Series', 'Rolling Drawdown Duration']:
        if col in combined_ts.columns:
            drawdown_data[col] = combined_ts[col]
    
    if not drawdown_data.empty:
        drawdown_file = os.path.join(output_dir, f"{prefix}_drawdown_{timestamp}.csv")
        drawdown_data.to_csv(drawdown_file)
        exported_files['drawdown'] = drawdown_file
    
    # Risk Metrics Over Time
    risk_data = pd.DataFrame()
    risk_cols = ['Rolling Sharpe', 'Rolling Sortino', 'Rolling Volatility', 
                 'Rolling Skew', 'Rolling Var']
    for col in risk_cols:
        if col in combined_ts.columns:
            risk_data[col] = combined_ts[col]
    
    if not risk_data.empty:
        risk_file = os.path.join(output_dir, f"{prefix}_risk_metrics_{timestamp}.csv")
        risk_data.to_csv(risk_file)
        exported_files['risk_metrics'] = risk_file


def _create_metadata_file(results, exported_files, output_dir, prefix, timestamp):
    """Helper function to create metadata file"""
    metadata = {
        'Export Timestamp': timestamp,
        'Files Exported': list(exported_files.keys()),
        'Start Date': results.get('start_date', 'N/A'),
        'End Date': results.get('end_date', 'N/A'),
        'Initial Capital': results.get('initial_capital', 'N/A'),
        'Final Capital': results.get('final_capital', 'N/A'),
        'Total Return': f"{results.get('returns', 0) * 100:.2f}%"
    }
    
    metadata_file = os.path.join(output_dir, f"{prefix}_metadata_{timestamp}.txt")
    with open(metadata_file, 'w') as f:
        for key, value in metadata.items():
            if key == 'Files Exported':
                f.write(f"{key}:\n")
                for file in value:
                    f.write(f"  - {file}\n")
            else:
                f.write(f"{key}: {value}\n")
    exported_files['metadata'] = metadata_file


import os
import re
import glob
from datetime import datetime
import pandas as pd
import matplotlib.pyplot as plt
from matplotlib.ticker import PercentFormatter
from pptx import Presentation
from pptx.util import Inches, Pt

# -----------------------------
# Helpers
# -----------------------------

def _pick_latest(path_pattern: str) -> str | None:
    """Return the latest file matching a pattern or None if none exist."""
    matches = glob.glob(path_pattern)
    if not matches:
        return None
    # Try to sort by timestamp embedded at *_YYYYmmdd_HHMMSS.*
    def key_fn(p):
        m = re.search(r'_(\d{8}_\d{6})\.', os.path.basename(p))
        if m:
            try:
                return datetime.strptime(m.group(1), "%Y%m%d_%H%M%S")
            except ValueError:
                pass
        return datetime.fromtimestamp(os.path.getmtime(p))
    return sorted(matches, key=key_fn)[-1]

def _load_timeseries_csv(path: str) -> pd.DataFrame:
    """Load a CSV with a 'Date' index or first column as datetime index."""
    df = pd.read_csv(path)
    if 'Date' in df.columns:
        df['Date'] = pd.to_datetime(df['Date'], errors='coerce')
        df = df.set_index('Date').sort_index()
    else:
        # heuristic: try first column
        first_col = df.columns[0]
        try:
            df[first_col] = pd.to_datetime(df[first_col], errors='coerce')
            if df[first_col].notna().any():
                df = df.set_index(first_col).sort_index()
        except Exception:
            pass
    return df

def _ensure_dir(d: str):
    os.makedirs(d, exist_ok=True)

def _save_line_chart(series_or_df, title, y_as_percent=False, outfile="chart.png"):
    plt.figure(figsize=(10, 5))
    if isinstance(series_or_df, pd.Series):
        series_or_df.plot()
    else:
        series_or_df.plot()
    plt.title(title)
    plt.xlabel("Date")
    plt.ylabel("Value")
    if y_as_percent:
        plt.gca().yaxis.set_major_formatter(PercentFormatter(1.0))
    plt.tight_layout()
    plt.savefig(outfile, dpi=200)
    plt.plot()
    plt.close()

    return outfile

def _pct_like(colname: str) -> bool:
    """Heuristic: treat columns with 'return', 'drawdown', 'relative', 'vol', 'sharpe', 'sortino', 'var' as percentages if they look like 0..1."""
    return any(k in colname.lower() for k in ["return", "drawdown", "relative", "vol", "sharpe", "sortino", "var"])

# -----------------------------
# Main entry
# -----------------------------

def create_charts_from_exports(
    output_dir="backtest_exports",
    prefix="backtest",
    charts_dir="charts",
    pptx_template: str | None = None,
    pptx_out: str | None = None
):
    """
    Scan export folder, create charts, and optionally assemble a PPTX.

    Parameters
    ----------
    output_dir : str
        Folder with your CSV exports.
    prefix : str
        Filename prefix used during export.
    charts_dir : str
        Output folder for generated PNG charts.
    pptx_template : str | None
        Optional PowerPoint template (e.g., 'template.pptx'). If None, creates a blank deck.
    pptx_out : str | None
        If provided, saves a PPTX with the generated charts embedded.
    """
    _ensure_dir(charts_dir)

    created = []   # list of (title, path)

    # 1) Combined timeseries (if present)
    ts_path = _pick_latest(os.path.join(output_dir, f"{prefix}_timeseries_*.csv"))
    if ts_path:
        ts_df = _load_timeseries_csv(ts_path)
        # Plot all columns (line chart)
        title = "Backtest Time Series"
        outfile = os.path.join(charts_dir, "timeseries_all.png")
        # Decide if percent-y is appropriate: only if *all* numeric columns look like fractions
        numeric_cols = ts_df.select_dtypes("number").columns
        y_as_percent = len(numeric_cols) > 0 and all(_pct_like(c) for c in numeric_cols)
        created.append((
            title,
            _save_line_chart(ts_df[numeric_cols], title, y_as_percent=y_as_percent, outfile=outfile)
        ))

        # Optional: create per-column charts (cleaner for slides)
        for col in numeric_cols:
            title = f"Time Series — {col}"
            outfile = os.path.join(charts_dir, f"timeseries_{re.sub(r'[^a-zA-Z0-9]+','_',col.lower())}.png")
            created.append((
                title,
                _save_line_chart(ts_df[col].dropna(), title, y_as_percent=_pct_like(col), outfile=outfile)
            ))

    # 2) Equity curves (Equity Curve, Benchmark Equity Curve, Relative Performance)
    eq_path = _pick_latest(os.path.join(output_dir, f"{prefix}_equity_curve_*.csv"))
    if eq_path:
        eq_df = _load_timeseries_csv(eq_path)
        # Combined equity curves
        cols = [c for c in eq_df.columns if "equity" in c.lower()]
        if cols:
            title = "Equity Curves"
            outfile = os.path.join(charts_dir, "equity_curves.png")
            created.append((
                title,
                _save_line_chart(eq_df[cols], title, y_as_percent=False, outfile=outfile)
            ))
        # Relative performance
        rel_cols = [c for c in eq_df.columns if "relative" in c.lower()]
        for col in rel_cols:
            title = f"Relative Performance — {col}"
            outfile = os.path.join(charts_dir, f"relative_{re.sub(r'[^a-zA-Z0-9]+','_',col.lower())}.png")
            created.append((
                title,
                _save_line_chart(eq_df[col].dropna(), title, y_as_percent=True, outfile=outfile)
            ))

    # 3) Drawdowns
    dd_path = _pick_latest(os.path.join(output_dir, f"{prefix}_drawdown_*.csv"))
    if dd_path:
        dd_df = _load_timeseries_csv(dd_path)
        # Plot each drawdown-like column individually
        for col in dd_df.columns:
            title = f"Drawdown — {col}"
            outfile = os.path.join(charts_dir, f"drawdown_{re.sub(r'[^a-zA-Z0-9]+','_',col.lower())}.png")
            created.append((
                title,
                _save_line_chart(dd_df[col].dropna(), title, y_as_percent=True, outfile=outfile)
            ))

    # 4) Risk metrics over time
    risk_path = _pick_latest(os.path.join(output_dir, f"{prefix}_risk_metrics_*.csv"))
    if risk_path:
        risk_df = _load_timeseries_csv(risk_path)
        for col in risk_df.columns:
            title = f"Risk Metric — {col}"
            outfile = os.path.join(charts_dir, f"risk_{re.sub(r'[^a-zA-Z0-9]+','_',col.lower())}.png")
            # Treat vol/sortino/sharpe/var heuristically as percent-like if values are in 0..1
            y_as_percent = _pct_like(col)
            created.append((
                title,
                _save_line_chart(risk_df[col].dropna(), title, y_as_percent=y_as_percent, outfile=outfile)
            ))

    # 5) Raw portfolio value (useful single line chart)
    pv_path = _pick_latest(os.path.join(output_dir, f"{prefix}_portfolio_value_*.csv"))
    if pv_path:
        pv_df = _load_timeseries_csv(pv_path)
        # try common column names
        value_col = None
        for c in pv_df.columns:
            if "value" in c.lower() or "equity" in c.lower() or "portfolio" in c.lower():
                value_col = c
                break
        if value_col:
            title = "Portfolio Value"
            outfile = os.path.join(charts_dir, "portfolio_value.png")
            created.append((
                title,
                _save_line_chart(pv_df[value_col].dropna(), title, y_as_percent=False, outfile=outfile)
            ))

    # -----------------------------
    # Optional: Build PPTX
    # -----------------------------
    if pptx_out:
        if pptx_template and os.path.exists(pptx_template):
            prs = Presentation(pptx_template)
        else:
            prs = Presentation()
            # add a title-only layout as fallback
            # (Index 5 is often Title Only; fall back to first if not available)
            title_only_layout = None
            for i, layout in enumerate(prs.slide_layouts):
                # heuristic: choose one with title placeholder
                if any(ph.placeholder_format.type == 1 for ph in layout.placeholders):  # 1 = TITLE
                    title_only_layout = layout
                    break
            if title_only_layout is None:
                title_only_layout = prs.slide_layouts[0]

        # Add one slide per chart
        for title, img_path in created:
            # pick a layout with title, otherwise first
            slide_layout = None
            for layout in prs.slide_layouts:
                if any(ph.placeholder_format.type == 1 for ph in layout.placeholders):
                    slide_layout = layout
                    break
            if slide_layout is None:
                slide_layout = prs.slide_layouts[0]

            slide = prs.slides.add_slide(slide_layout)
            # Set title (if placeholder exists)
            if slide.shapes.title:
                slide.shapes.title.text = title

            # Add image centered with margins
            pic_left = Inches(0.7)
            pic_top = Inches(1.6)
            pic_height = Inches(5.0)
            slide.shapes.add_picture(img_path, pic_left, pic_top, height=pic_height)

        prs.save(pptx_out)

    return created

# -----------------------------
# Example usage (uncomment):
# -----------------------------
# charts = create_charts_from_exports(
#     output_dir="backtest_exports",
#     prefix="backtest",
#     charts_dir="charts",
#     pptx_template="template.pptx",   # or None
#     pptx_out="backtest_charts.pptx"  # or None to skip pptx creation
# )
# print("Created charts:")
# for title, path in charts:
#     print(f"- {title}: {path}")
